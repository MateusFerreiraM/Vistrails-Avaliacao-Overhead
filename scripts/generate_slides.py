from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide_with_image(prs, title, img_path):
    slide_layout = prs.slide_layouts[5] # Blank slide with title
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Image
    try:
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(img_path, left, top, height=height)
    except Exception as e:
        print(f"Erro ao adicionar imagem {img_path}: {e}")
        
    return slide

def main():
    prs = Presentation()
    
    # 1. Slide de Título
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Comparação de Desempenho:\nPython Puro vs VisTrailsJL"
    subtitle.text = "Análise de Overhead e Benefícios da Proveniência"
    
    # 2. Slide de Metodologia
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Metodologia: 10 Workflows"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Mapeamos 10 fluxos de dados distintos para cobrir diversas complexidades:"
    p = tf.add_paragraph()
    p.text = "1. Numéricos Leves (gcd, primes, outputs)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Renderização de Gráficos (scatter, bar_ex1, hist_ex1, lineplot_ex3)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Treinamento de Machine Learning (pipeline, grid_search)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Processamento de Imagens (imagemagick)"
    p.level = 1
    
    # 3. Slide Gráfico Tempo Médio
    add_slide_with_image(prs, "Resultados: Tempo Médio de Execução", "resultados/grafico_tempo_medio.png")
    
    # 4. Slide Gráfico Boxplot
    add_slide_with_image(prs, "Resultados: Variação e Estabilidade (Boxplot)", "resultados/grafico_boxplot.png")

    # 5. Slide Gráfico Logarítmico
    add_slide_with_image(prs, "Resultados: Escala Logarítmica (Detalhe de Tarefas Leves)", "resultados/grafico_tempo_log.png")
    
    # 6. Slide Gráfico Overhead
    add_slide_with_image(prs, "Análise do Overhead Injetado pelo VisTrails", "resultados/grafico_overhead.png")
    
    # 7. Slide Conclusão
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusões"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "O Overhead não é linear:"
    
    p = tf.add_paragraph()
    p.text = "Tarefas extremamente rápidas (Python = 0.1s) sofrem muito (VisTrails = 3.0s). Ocorre penalidade pela Injeção PyCall e Action Replay."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Em fluxos que realmente demandam processamento (Treinamento de ML, Renderização Gráfica), o VisTrails apenas adiciona 1 a 3 segundos de overhead constante."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Conclusão Final: Para Ciência de Dados real, o VisTrails é viável. A proveniência, auditoria de estados passados e reprodutibilidade compensam totalmente os poucos segundos de overhead de compilação gráfica inicial."
    p.level = 1
    
    # Save presentation
    prs.save("Apresentacao_VisTrails.pptx")
    print("Apresentação gerada com sucesso: Apresentacao_VisTrails.pptx")

if __name__ == '__main__':
    main()
